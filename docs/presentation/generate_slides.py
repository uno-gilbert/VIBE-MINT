#!/usr/bin/env python3
"""Generate VibeMint workshop PowerPoint with infographic images."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
IMAGES = ROOT / "images"
OUTPUT = ROOT / "VibeMint-Workshop-2026.pptx"

WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(15, 23, 42)


def set_slide_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 15, 23, 42)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(191, 219, 254)


def add_image_slide(prs, image_name: str):
    path = IMAGES / image_name
    if not path.exists():
        raise FileNotFoundError(f"Missing slide image: {path}")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 255, 255, 255)
    slide.shapes.add_picture(
        str(path),
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )


def add_lunch_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 37, 99, 235)
    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1.0))
    t1.text_frame.paragraphs[0].text = "12:30 ~ 13:30"
    t1.text_frame.paragraphs[0].font.size = Pt(28)
    t1.text_frame.paragraphs[0].font.color.rgb = RGBColor(219, 234, 254)
    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11), Inches(1.5))
    t2.text_frame.paragraphs[0].text = "점심 식사"
    t2.text_frame.paragraphs[0].font.size = Pt(44)
    t2.text_frame.paragraphs[0].font.bold = True
    t2.text_frame.paragraphs[0].font.color.rgb = WHITE
    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(1.0))
    t3.text_frame.paragraphs[0].text = "점심 후: MetaMask Sepolia · Cursor 로그인 재확인"
    t3.text_frame.paragraphs[0].font.size = Pt(20)
    t3.text_frame.paragraphs[0].font.color.rgb = RGBColor(219, 234, 254)


# (중제목 PNG, [소제목 PNG ...])
SLIDE_GROUPS = [
    (
        "session1-1-orientation.png",
        [
            "session1-1-1-instructor-env.png",
            "session1-1-2-vibe-workflow.png",
            "session1-1-3-prep-checklist.png",
        ],
    ),
    (
        "session1-2-nft-erc.png",
        [
            "session1-2-1-nft-concepts.png",
            "session1-2-2-erc-openzeppelin.png",
            "session1-2-3-cursor-reverse.png",
        ],
    ),
    (
        "session1-3-spec.png",
        [
            "session1-3-1-nft-benchmark.png",
            "session1-3-2-spec-template.png",
            "session1-3-3-spec-confirm.png",
        ],
    ),
    (
        "session2-1-remix-rules.png",
        [
            "session2-1-1-00-rules.png",
            "session2-1-2-remix-vm.png",
            "session2-1-3-compile-deploy.png",
        ],
    ),
    (
        "session2-2-stage-build.png",
        [
            "session2-2-1-stage-0-1.png",
            "session2-2-2-stage-2.png",
            "session2-2-3-stage-3.png",
        ],
    ),
    (
        "session2-3-audit.png",
        [
            "session2-3-1-audit-run.png",
            "session2-3-2-critical-fix.png",
            "session2-3-3-retest.png",
        ],
    ),
    (
        "session3-1-sepolia-deploy.png",
        [
            "session3-1-1-faucet-deploy.png",
            "session3-1-2-seturi-mint.png",
            "session3-1-3-etherscan-multichain.png",
        ],
    ),
    (
        "session3-2-frontend-dapp.png",
        [
            "session3-2-1-env-setup.png",
            "session3-2-2-npm-dev.png",
            "session3-2-3-connect-mint.png",
        ],
    ),
    (
        "session3-3-opensea-wrap.png",
        [
            "session3-3-1-opensea.png",
            "session3-3-2-ui-polish.png",
            "session3-3-3-wrap-qa.png",
        ],
    ),
]


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "쉽게 이해하고 빠르게\nAI 바이브 코딩으로 만들어 보는 NFT",
        "VibeMint 워크숍 · 2026 · Sepolia 테스트넷",
    )

    for i, (section_image, subtitles) in enumerate(SLIDE_GROUPS):
        if i == 3:
            add_lunch_slide(prs)
        add_image_slide(prs, section_image)
        for subtitle_image in subtitles:
            add_image_slide(prs, subtitle_image)

    add_title_slide(prs, "감사합니다", "Q & A · Sepolia 교육용 — 메인넷 전문 Audit 필수")

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
